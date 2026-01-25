import os
import sys
import subprocess
from pathlib import Path

def install_from_requirements():
    """requirements.txt fayldan kutubxonalarni o'rnatadi"""
    req_file = Path("requirements.txt")
    
    if req_file.exists():
        print("requirements.txt topildi. Kutubxonalar o'rnatilmoqda...")
        try:
            subprocess.check_call([sys.executable, "-m", "pip", "install", "-r", "requirements.txt"])
            print("✓ Barcha kutubxonalar o'rnatildi")
            return True
        except subprocess.CalledProcessError as e:
            print(f"✗ O'rnatishda xatolik: {e}")
            return False
    else:
        print("requirements.txt topilmadi")
        return False

def check_requirements():
    """Kerakli kutubxonalar mavjudligini tekshiradi"""
    required_modules = ['docx2pdf', 'docx', 'reportlab', 'PIL', 'pandas']
    missing = []
    
    for module in required_modules:
        try:
            __import__(module)
        except ImportError:
            missing.append(module)
    
    if missing:
        print(f"Kerakli kutubxonalar topilmadi: {', '.join(missing)}")
        choice = input("requirements.txt dan o'rnatishni xohlaysizmi? (y/n): ")
        if choice.lower() == 'y':
            return install_from_requirements()
    return True

def convert_any_to_pdf(input_path, output_path=None):
    """
    Har qanday faylni asl formatini saqlagan holda PDFga aylantiradi
    """
    input_path = Path(input_path)
    if not input_path.exists():
        raise FileNotFoundError(f"Fayl topilmadi: {input_path}")
    
    if output_path is None:
        output_path = input_path.with_suffix('.pdf')
    else:
        output_path = Path(output_path)
    
    ext = input_path.suffix.lower()
    
    try:
        if ext in ['.docx', '.doc']:
            convert_word_to_pdf(input_path, output_path)
        elif ext in ['.pptx', '.ppt']:
            convert_powerpoint_to_pdf(input_path, output_path)
        elif ext in ['.xlsx', '.xls']:
            convert_excel_to_pdf(input_path, output_path)
        elif ext in ['.html', '.htm']:
            convert_html_to_pdf(input_path, output_path)
        elif ext in ['.txt']:
            convert_text_to_pdf(input_path, output_path)
        elif ext in ['.jpg', '.jpeg', '.png', '.bmp', '.gif']:
            convert_image_to_pdf(input_path, output_path)
        elif ext == '.pdf':
            print("Fayl allaqachon PDF formatida!")
            return str(input_path)
        else:
            print(f"Bu format qo'llab-quvvatlanmaydi: {ext}")
            return None
            
        print(f"Muvaffaqiyatli aylantildi: {output_path}")
        return str(output_path)
        
    except ImportError as e:
        print(f"Kerakli kutubxona o'rnatilmagan: {e}")
        print("Quyidagi buyruqni bajaring: pip install -r requirements.txt")
        return None
    except Exception as e:
        print(f"Xatolik: {e}")
        return None

def convert_word_to_pdf(input_path, output_path):
    """Word fayllarini PDFga aylantiradi - asl formatni saqlaydi"""
    try:
        # docx2pdf - eng yaxshi variant, asl formatni to'liq saqlaydi
        from docx2pdf import convert
        convert(str(input_path), str(output_path))
    except ImportError:
        # python-docx va reportlab - formatni yaxshiroq saqlash
        import docx
        from reportlab.pdfgen import canvas
        from reportlab.lib.pagesizes import letter
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
        
        doc = docx.Document(str(input_path))
        pdf_doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                # Bold va italic formatni saqlash
                style = styles['Normal']
                if paragraph.runs:
                    for run in paragraph.runs:
                        if run.bold:
                            style = styles['Heading3']
                        elif run.italic:
                            style = styles['Italic']
                
                p = Paragraph(paragraph.text, style)
                story.append(p)
                story.append(Spacer(1, 12))
        
        pdf_doc.build(story)

def convert_powerpoint_to_pdf(input_path, output_path):
    """PowerPoint fayllarini PDFga aylantiradi - slaydlarni saqlaydi"""
    try:
        # comtypes - eng yaxshi variant (Windows)
        import comtypes.client
        powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
        powerpoint.Visible = 1
        deck = powerpoint.Presentations.Open(str(input_path.absolute()))
        deck.SaveAs(str(output_path.absolute()), 32)  # 32 = PDF format
        deck.Close()
        powerpoint.Quit()
    except:
        # python-pptx va reportlab
        from pptx import Presentation
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib.pagesizes import letter
        
        prs = Presentation(str(input_path))
        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        for i, slide in enumerate(prs.slides):
            story.append(Paragraph(f"Slayd {i+1}", styles['Heading1']))
            story.append(Spacer(1, 20))
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    p = Paragraph(shape.text, styles['Normal'])
                    story.append(p)
                    story.append(Spacer(1, 12))
            
            if i < len(prs.slides) - 1:
                story.append(PageBreak())
        
        doc.build(story)

def convert_excel_to_pdf(input_path, output_path):
    """Excel fayllarini PDFga aylantiradi - jadval formatini saqlaydi"""
    import pandas as pd
    from reportlab.lib.pagesizes import letter, A4
    from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
    from reportlab.lib import colors
    
    # Excel faylni o'qish
    df = pd.read_excel(str(input_path))
    
    # PDF yaratish
    doc = SimpleDocTemplate(str(output_path), pagesize=A4)
    elements = []
    
    # Ma'lumotlarni jadval formatida tayyorlash
    data = [df.columns.tolist()]  # Header
    for index, row in df.iterrows():
        data.append(row.tolist())
    
    # Jadval yaratish
    table = Table(data)
    
    # Jadval stilini belgilash
    table.setStyle(TableStyle([
        ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
        ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
        ('ALIGN', (0, 0), (-1, -1), 'CENTER'),
        ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
        ('FONTSIZE', (0, 0), (-1, 0), 12),
        ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
        ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
        ('GRID', (0, 0), (-1, -1), 1, colors.black)
    ]))
    
    elements.append(table)
    doc.build(elements)

def convert_html_to_pdf(input_path, output_path):
    """HTML fayllarini PDFga aylantiradi"""
    try:
        import pdfkit
        pdfkit.from_file(str(input_path), str(output_path))
    except ImportError:
        try:
            from weasyprint import HTML
            HTML(filename=str(input_path)).write_pdf(str(output_path))
        except ImportError:
            print("pdfkit yoki weasyprint o'rnating")

def convert_text_to_pdf(input_path, output_path):
    """Text fayllarini PDFga aylantiradi"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    
    c = canvas.Canvas(str(output_path), pagesize=letter)
    
    with open(input_path, 'r', encoding='utf-8') as file:
        lines = file.readlines()
    
    y = 750
    for line in lines:
        if y < 50:
            c.showPage()
            y = 750
        c.drawString(50, y, line.strip())
        y -= 15
    c.save()

def convert_image_to_pdf(input_path, output_path):
    """Rasm fayllarini PDFga aylantiradi - asl sifatni saqlaydi"""
    from reportlab.pdfgen import canvas
    from reportlab.lib.pagesizes import letter
    from PIL import Image
    
    img = Image.open(str(input_path))
    img_width, img_height = img.size
    
    # PDF o'lchamini rasmga moslashtirish
    if img_width > img_height:
        # Landscape
        from reportlab.lib.pagesizes import landscape
        pagesize = landscape(letter)
    else:
        # Portrait
        pagesize = letter
    
    c = canvas.Canvas(str(output_path), pagesize=pagesize)
    page_width, page_height = pagesize
    
    # Rasmni to'liq sahifaga moslashtirish
    scale = min(page_width/img_width, page_height/img_height) * 0.95
    new_width = img_width * scale
    new_height = img_height * scale
    
    x = (page_width - new_width) / 2
    y = (page_height - new_height) / 2
    
    # Yuqori sifatda rasm qo'shish
    c.drawImage(str(input_path), x, y, width=new_width, height=new_height, preserveAspectRatio=True)
    c.save()

def batch_convert(folder_path, output_folder=None):
    """Papkadagi barcha fayllarni PDFga aylantiradi"""
    folder = Path(folder_path)
    if output_folder is None:
        output_folder = folder / "pdf_output"
    else:
        output_folder = Path(output_folder)
    
    output_folder.mkdir(exist_ok=True)
    
    converted_files = []
    for file_path in folder.iterdir():
        if file_path.is_file() and file_path.suffix.lower() != '.pdf':
            output_path = output_folder / f"{file_path.stem}.pdf"
            result = convert_any_to_pdf(file_path, output_path)
            if result:
                converted_files.append(result)
    
    print(f"{len(converted_files)} ta fayl aylantildi")
    return converted_files

if __name__ == "__main__":
    # Requirements tekshirish
    if not check_requirements():
        print("Kerakli kutubxonalar o'rnatilmadi. Dastur to'xtatildi.")
        sys.exit(1)
    
    if len(sys.argv) < 2:
        print("Foydalanish:")
        print("1. Bitta fayl: python universal_pdf_converter.py file.docx") 
        print("2. Output belgilash: python universal_pdf_converter.py file.docx output.pdf")
        print("3. Papka: python universal_pdf_converter.py --folder ./documents")
        sys.exit(1)
    
    if sys.argv[1] == "--folder":
        folder_path = sys.argv[2] if len(sys.argv) > 2 else "."
        batch_convert(folder_path)
    else:
        input_file = sys.argv[1]
        output_file = sys.argv[2] if len(sys.argv) > 2 else None
        convert_any_to_pdf(input_file, output_file)
